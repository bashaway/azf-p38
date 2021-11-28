import sys
import io
import os
from pptx import Presentation
import logging
import http.client, urllib.request, urllib.parse, urllib.error, base64, json
import azure.functions as func
from base64 import b64decode
import fitz

# Image は PDFが読み込めない
#from PIL import Image

# pdf2image は popplerのエラーがでる
#from pdf2image import convert_from_path, convert_from_bytes

# wandはimportするだけでエラーになる
#from wand.image import Image

def main(req: func.HttpRequest) -> func.HttpResponse:
    logging.info('Python HTTP trigger function processed a request.')

    try:
        req_body = req.get_json()
    except ValueError:
        dict_return = {
            "string": "no input file : print sample JSON schema",
            "pages": 2,
            "results": [{"page":1, "note": "string:PAGE1 note body", "image": "base64: PAGE1 image body"},
                        {"page":2, "note": "string:PAGE2 note body", "image": "base64: PAGE2 image body"}
                       ]
        }

        return func.HttpResponse(json.dumps(dict_return, sort_keys=True, indent=2) , mimetype="application/json" )

    req_body = req.get_json()

    body_pptx = req_body.get('pptx')
    bin_pptx = b64decode(body_pptx)
    byte_pptx = io.BytesIO( bin_pptx )

    presentation = Presentation( byte_pptx )

    list_note = []
    list_info = []
    for i, slide in enumerate(presentation.slides):
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text:
            list_note.append({'page':i+1,'note':slide.notes_slide.notes_text_frame.text})
            list_info.append({'page':i+1,'note':base64.b64encode(slide.notes_slide.notes_text_frame.text.encode()).decode('utf-8')})
            #list_info.append({'page':i+1,'note':slide.notes_slide.notes_text_frame.text})


    body_pdf = req_body.get('pdf')
    bin_pdf = b64decode(body_pdf)
    byte_pdf = io.BytesIO( bin_pdf )

    # fitz make success
    doc = fitz.open(stream=bin_pdf, filetype='pdf')

    print(len(doc))
    image_matrix = fitz.Matrix(fitz.Identity)
    image_matrix.preScale(2, 2)

    list_image = []
    for i, page in enumerate(doc):
        pix = page.getPixmap(alpha = False, matrix=image_matrix)
        body = pix.getImageData(output="png")
        str_base64encoded = base64.b64encode(body).decode('utf-8')
        list_image.append({'page':i+1,'image':str_base64encoded})

        for x in list_info:
            if x['page'] == (i+1):
                x['image'] = str_base64encoded


    dict_return = {
        "string": "String",
        "pages": i,
        "results": list_info
    }

    return func.HttpResponse(json.dumps(dict_return, sort_keys=True, indent=2) , mimetype="application/json" )